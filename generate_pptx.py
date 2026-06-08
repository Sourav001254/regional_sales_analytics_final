from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import sys

# Add current directory to path to import src.logger
sys.path.append(os.path.dirname(os.path.abspath(__file__)))
from logger import logger

def create_executive_summary():
    logger.info("="*50)
    logger.info("STAGE 7a: AUTO-GENERATING POWERPOINT SUMMARY")
    logger.info("="*50)
    
    prs = Presentation()
    
    # Define a clean layout
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Regional Sales Performance Analysis"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(10, 22, 43) # Navy
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "Executive Summary & Advanced Insights\nGenerated Automatically via Python Pipeline"
    
    # Slide 2: Executive Summary
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    
    title2 = slide2.shapes.title
    title2.text = "Key Business Insights"
    
    body_shape = slide2.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Automated Data Findings:"
    
    # Load insights if available
    insights_path = "advanced_outputs/automated_insights.txt"
    if os.path.exists(insights_path):
        with open(insights_path, "r") as f:
            lines = f.readlines()
            for line in lines:
                p = tf.add_paragraph()
                p.text = line.strip()
                p.level = 1
    else:
        p = tf.add_paragraph()
        p.text = "Insights file not found. Ensure Step 6 was executed."
        p.level = 1
        
    # Attempt to insert an image if EDA was run
    # Adding a new slide for visual
    blank_slide_layout = prs.slide_layouts[5]
    slide3 = prs.slides.add_slide(blank_slide_layout)
    slide3.shapes.title.text = "Revenue Forecast (SARIMA)"
    
    img_path = "advanced_outputs/revenue_forecast_sarima.png"
    if os.path.exists(img_path):
        left = Inches(1)
        top = Inches(2)
        height = Inches(4.5)
        slide3.shapes.add_picture(img_path, left, top, height=height)
    else:
        txBox = slide3.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        txBox.text_frame.text = "Forecast Chart Image not found (Run Step 6)."

    # Slide 4: Regional performance bar chart
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    slide4.shapes.title.text = "Regional Performance Snapshot"
    
    # Slide 5: RFM segments
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    slide5.shapes.title.text = "Customer Cohorts & RFM Segments"
    
    # Slide 6: Thank You
    slide6 = prs.slides.add_slide(prs.slide_layouts[0])
    slide6.shapes.title.text = "Thank You"
    slide6.placeholders[1].text = "Contact: analytics-team@company.com"

    # Save
    output_filename = "Executive_Summary_Report.pptx"
    prs.save(output_filename)
    logger.info(f"[v] Successfully generated PowerPoint: {output_filename}")

if __name__ == "__main__":
    try:
        import pptx
    except ImportError:
        logger.error("Note: 'python-pptx' is required to run this script. Please pip install it first.")
    
    create_executive_summary()
